from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Tạo PowerPoint
prs = Presentation()

# Hàm thêm slide với tiêu đề + nội dung
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(14)
    p.font.name = 'Arial'
    return slide

# Nội dung các slide từ vựng
slides_data = [
    ("Thời gian (Time)", "幾 (jǐ) – N – How many?\n\n點/點鐘 (diǎn/diǎnzhōng) – M – o’clock\n例句: 現在幾點鐘？\nXiànzài jǐ diǎnzhōng?\n\n分 (fēn) – M – minute"),
    ("Ngày tháng (Date)", "今天 (jīntiān) – N – today\n星期 (xīngqí) – N – week; day of the week\n月 (yuè) – N – month\n號 (hào) – M – day of month\n例句: 今天是五月十七日, 星期一。\nJīntiān shì wǔ yuè shíqī rì, xīngqí yī.\n日 (rì) – N – day\n星期天 (xīngqítiān) – N – Sunday"),
    ("Buổi trong ngày (Parts of the day)", "早上 (zǎoshàng) – N – morning\n上午 (shàngwǔ) – N – morning (before noon)\n中午 (zhōngwǔ) – N – noon\n下午 (xiàwǔ) – N – afternoon"),
    ("Trường học & Học tập", "學校 (xuéxiào) – N – school\n例句: 你幾點去學校？\nNǐ jǐ diǎn qù xuéxiào?\n課 (kè) – N – class\n上課 (shàngkè) – V-sep – to attend class\n下課 (xiàkè) – V-sep – class dismissed\n圖書館 (túshūguǎn) – N – library"),
    ("Gia đình (Family)", "媽媽 (māma) – N – mom\n爸爸 (bàba) – N – dad\n孩子 (háizi) – N – child; children\n家 (jiā) – N – home; house\n家人 (jiārén) – N – family members"),
    ("Động từ thường dùng (Common Verbs)", "去 (qù) – V – to go\n回 (huí) – V – to return\n有 (yǒu) – Vst – to have\n沒 (méi) – Adv – not")
]

# Thêm slide từ vựng
for title, content in slides_data:
    add_slide(title, content)

# Luyện tập
practice_data = [
    ("Luyện tập 1 – Hỏi giờ", 
     "A: 現在幾點？ (Xiànzài jǐ diǎn?)\nB: 現在十點二十五分。 (Xiànzài shí diǎn èrshíwǔ fēn.)"),
    ("Luyện tập 2 – Ngày tháng", 
     "A: 今天星期幾？ (Jīntiān xīngqí jǐ?)\nB: 今天星期三。 (Jīntiān xīngqí sān.)\n\nA: 今天幾月幾號？ (Jīntiān jǐ yuè jǐ hào?)\nB: 今天十二月七號。 (Jīntiān shí’èr yuè qī hào.)"),
    ("Luyện tập 3 – Thời gian + Hành động", 
     "A: 你幾點去學校？ (Nǐ jǐ diǎn qù xuéxiào?)\nB: 我早上七點去學校。 (Wǒ zǎoshàng qī diǎn qù xuéxiào.)\n\nA: 你幾點下課？ (Nǐ jǐ diǎn xiàkè?)\nB: 我下午四點四十分下課。 (Wǒ xiàwǔ sì diǎn sìshí fēn xiàkè.)"),
    ("Luyện tập 4 – Có / Không có", 
     "A: 你今天上午有課嗎？ (Nǐ jīntiān shàngwǔ yǒu kè ma?)\nB: 我上午有課。 / 我上午沒有課。\n(Wǒ shàngwǔ yǒu kè. / Wǒ shàngwǔ méiyǒu kè.)\n\nA: 誰有圖書館的書？ (Shéi yǒu túshūguǎn de shū?)\nB: 我有圖書館的書。 (Wǒ yǒu túshūguǎn de shū.)")
]

for title, content in practice_data:
    add_slide(title, content)

# Đối thoại cuối
dialogue_content = """(在中明家 At Zhongming's place)

媽媽: 今天是星期一, 你不去學校嗎？
Māma: Jīntiān shì xīngqí yī, nǐ bù qù xuéxiào ma?

中明: 我早上沒有課, 九點去圖書館。
Zhōngmíng: Wǒ zǎoshàng méiyǒu kè, jiǔ diǎn qù túshūguǎn.

媽媽: 中午呢? 你中午回家嗎？
Māma: Zhōngwǔ ne? Nǐ zhōngwǔ huí jiā ma?

中明: 我下午有課, 中午不回家。
Zhōngmíng: Wǒ xiàwǔ yǒu kè, zhōngwǔ bù huí jiā.

媽媽: 你幾點下課?
Māma: Nǐ jǐ diǎn xiàkè?

中明: 我四點四十分下課。
Zhōngmíng: Wǒ sì diǎn sìshí fēn xiàkè.
"""
add_slide("對話一 – Dialogue 1", dialogue_content)

# Lưu file
pptx_path = "/mnt/data/Bai_2_Bai_Gioi_Thieu.pptx"
prs.save(pptx_path)
pptx_path